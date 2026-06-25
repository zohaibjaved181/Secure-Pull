#!/usr/bin/env python3
"""Generate Secure Pull Architecture PPTX presentation."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import copy

# Create presentation (widescreen 16:9)
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Color palette
RED = RGBColor(0xDC, 0x35, 0x45)
DARK_RED = RGBColor(0xA7, 0x1D, 0x2A)
ORANGE = RGBColor(0xFD, 0x7E, 0x14)
CYAN = RGBColor(0x00, 0xBC, 0xEB)
BLUE = RGBColor(0x00, 0x76, 0xD5)
GREEN = RGBColor(0x6C, 0xC0, 0x4A)
DARK_GREEN = RGBColor(0x19, 0x87, 0x54)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK_BG = RGBColor(0x0D, 0x1B, 0x2A)
LIGHT_GRAY = RGBColor(0xE8, 0xE8, 0xE8)

def add_rounded_rect(slide, left, top, width, height, fill_color, border_color=None, text="", font_size=10, font_color=WHITE, bold=False):
    """Add a rounded rectangle shape."""
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1.5)
    else:
        shape.line.fill.background()
    if text:
        tf = shape.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = text
        run.font.size = Pt(font_size)
        run.font.color.rgb = font_color
        run.font.bold = bold
    return shape

def add_circle(slide, left, top, diameter, fill_color, text="", font_size=10, font_color=WHITE, bold=False):
    """Add a circle/oval shape."""
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, diameter, diameter)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    if text:
        tf = shape.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = text
        run.font.size = Pt(font_size)
        run.font.color.rgb = font_color
        run.font.bold = bold
    return shape

def add_text_box(slide, left, top, width, height, text, font_size=10, font_color=WHITE, bold=False, alignment=PP_ALIGN.LEFT):
    """Add a text box."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = alignment
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.color.rgb = font_color
    run.font.bold = bold
    return txBox

def add_arrow(slide, start_left, start_top, end_left, end_top, color):
    """Add a connector/arrow line."""
    connector = slide.shapes.add_connector(1, start_left, start_top, end_left, end_top)
    connector.line.color.rgb = color
    connector.line.width = Pt(2)
    return connector

def add_layer_row(slide, y_pos, layer_num, badge_color, title, details, row_color, row_border):
    """Add a complete layer row to the slide."""
    # Row background
    row_bg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.3), y_pos, Inches(12.7), Inches(1.15))
    row_bg.fill.solid()
    row_bg.fill.fore_color.rgb = row_color
    row_bg.line.color.rgb = row_border
    row_bg.line.width = Pt(1)

    # Layer badge (circle)
    badge = add_circle(slide, Inches(0.5), y_pos + Inches(0.15), Inches(0.85), badge_color,
                       f"LAYER\n{layer_num}", font_size=9, bold=True)

    # Layer title
    add_text_box(slide, Inches(1.5), y_pos + Inches(0.2), Inches(1.5), Inches(0.8),
                 title, font_size=14, font_color=WHITE, bold=True)

    # Firewall icon
    fw_left = Inches(3.2)
    fw = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, fw_left, y_pos + Inches(0.25), Inches(0.6), Inches(0.65))
    fw.fill.solid()
    fw.fill.fore_color.rgb = RED if layer_num >= 4 else (CYAN if layer_num == 3 else GREEN)
    fw.line.fill.background()

    # Switch icon
    sw_left = Inches(5.0)
    sw = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, sw_left, y_pos + Inches(0.35), Inches(0.8), Inches(0.35))
    sw.fill.solid()
    sw.fill.fore_color.rgb = BLUE
    sw.line.fill.background()

    # Server/VM icon
    srv_left = Inches(6.5)
    srv = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, srv_left, y_pos + Inches(0.25), Inches(0.5), Inches(0.65))
    srv.fill.solid()
    srv.fill.fore_color.rgb = CYAN
    srv.line.color.rgb = BLUE
    srv.line.width = Pt(1)

    # Traffic flow arrows
    arrow1 = add_arrow(slide, Inches(3.9), y_pos + Inches(0.55), Inches(4.9), y_pos + Inches(0.55), badge_color)
    arrow2 = add_arrow(slide, Inches(5.9), y_pos + Inches(0.55), Inches(6.4), y_pos + Inches(0.55), badge_color)

    # Flow indicator dots
    for i, offset in enumerate([Inches(4.2), Inches(4.5), Inches(4.7)]):
        dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, offset, y_pos + Inches(0.48), Inches(0.12), Inches(0.12))
        dot.fill.solid()
        dot.fill.fore_color.rgb = badge_color
        dot.line.fill.background()

    # Details bullet points
    detail_left = Inches(7.5)
    for i, detail in enumerate(details):
        add_text_box(slide, detail_left, y_pos + Inches(0.05 + i * 0.22), Inches(5.3), Inches(0.25),
                     f"\u2022 {detail}", font_size=8, font_color=LIGHT_GRAY)

# ============================================
# SLIDE 1: Title Slide
# ============================================
slide_layout = prs.slide_layouts[6]  # Blank
slide1 = prs.slides.add_slide(slide_layout)

# Dark background
bg = slide1.background
fill = bg.fill
fill.solid()
fill.fore_color.rgb = DARK_BG

# Title banner
title_shape = add_rounded_rect(slide1, Inches(2.5), Inches(2.5), Inches(8.3), Inches(1.2),
                                CYAN, text="Secure Pull Architecture", font_size=36, bold=True)

# Subtitle
add_text_box(slide1, Inches(3), Inches(4.2), Inches(7.3), Inches(0.8),
             "Multi-Layer Security Architecture with Unidirectional Traffic Flow",
             font_size=16, font_color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

# Decorative elements
add_rounded_rect(slide1, Inches(0.5), Inches(0.5), Inches(1.5), Inches(0.5),
                 RGBColor(0x1B, 0x28, 0x38), border_color=CYAN, text="SECURE PULL", font_size=9, bold=True)

# Zone labels
add_rounded_rect(slide1, Inches(10.5), Inches(5.5), Inches(1.2), Inches(0.4),
                 CYAN, text="DMZ", font_size=10, bold=True)
add_rounded_rect(slide1, Inches(10.5), Inches(6.1), Inches(1.2), Inches(0.4),
                 GREEN, text="DC", font_size=10, bold=True)

# ============================================
# SLIDE 2: Full Architecture Diagram
# ============================================
slide2 = prs.slides.add_slide(slide_layout)
bg2 = slide2.background
fill2 = bg2.fill
fill2.solid()
fill2.fore_color.rgb = DARK_BG

# Title
add_rounded_rect(slide2, Inches(3.5), Inches(0.15), Inches(6.3), Inches(0.55),
                 CYAN, text="Secure Pull Architecture", font_size=18, bold=True)

# Zone badges
add_rounded_rect(slide2, Inches(11.8), Inches(0.15), Inches(1.0), Inches(0.4),
                 CYAN, text="DMZ", font_size=9, bold=True)
add_rounded_rect(slide2, Inches(11.8), Inches(6.9), Inches(1.0), Inches(0.4),
                 GREEN, text="DC", font_size=9, bold=True)

# Secure Pull badge
add_rounded_rect(slide2, Inches(0.2), Inches(0.15), Inches(1.5), Inches(0.4),
                 RGBColor(0x1B, 0x28, 0x38), border_color=CYAN, text="SECURE PULL", font_size=8, bold=True)

# Layer 5: Internet Firewall
y = Inches(0.85)
layer5_details = [
    "Source IP Whitelisting - Only from Web Proxy",
    "Destination Server IP & URL Whitelisting",
    "IPS, Web Filtering & Anti-malware Security",
    "Zoning & Multi-tenancy",
    "Unidirectional Flow"
]
add_layer_row(slide2, y, 5, RED, "Internet\nFirewall", layer5_details,
              RGBColor(0x2D, 0x15, 0x1A), RGBColor(0x6E, 0x1A, 0x23))

# Vertical connector arrow 5->4
add_arrow(slide2, Inches(4.5), Inches(2.05), Inches(4.5), Inches(2.3), ORANGE)
add_text_box(slide2, Inches(4.7), Inches(2.05), Inches(0.8), Inches(0.25),
             "Proxy", font_size=7, font_color=RGBColor(0x99, 0x99, 0x99))

# Layer 4: Web Proxy
y = Inches(2.3)
layer4_details = [
    "Source IP Whitelisting - Only from Specific DMZ Servers",
    "Destination Server IP & URL Whitelisting",
    "Web Reputation & Web Usage Control",
    "Anti-malware Analysis & HTTPS Inspection"
]
add_layer_row(slide2, y, 4, ORANGE, "Web\nProxy", layer4_details,
              RGBColor(0x2D, 0x20, 0x0A), RGBColor(0x7F, 0x3F, 0x0A))

# Left note
add_text_box(slide2, Inches(0.15), Inches(2.5), Inches(1.4), Inches(0.6),
             "Only Proxy can communicate unidirectionally with Internet",
             font_size=6, font_color=RGBColor(0x88, 0xBB, 0xDD))

# Vertical connector 4->3
add_arrow(slide2, Inches(4.5), Inches(3.5), Inches(4.5), Inches(3.75), CYAN)
add_text_box(slide2, Inches(4.7), Inches(3.5), Inches(0.8), Inches(0.25),
             "Proxy", font_size=7, font_color=RGBColor(0x99, 0x99, 0x99))

# Layer 3: DMZ Firewall
y = Inches(3.75)
layer3_details = [
    "Source IP Whitelisting - Only from Specific DMZ Servers",
    "Destination Server IP & URL Whitelisting",
    "IPS, Web Filtering & Anti-malware Security",
    "Zoning & Unidirectional Flow"
]
add_layer_row(slide2, y, 3, CYAN, "DMZ\nFirewall", layer3_details,
              RGBColor(0x0A, 0x20, 0x2D), RGBColor(0x00, 0x5E, 0x76))

# Left note
add_text_box(slide2, Inches(0.15), Inches(3.95), Inches(1.4), Inches(0.6),
             "Only specific DMZ servers can communicate unidirectionally with Proxy",
             font_size=6, font_color=RGBColor(0x88, 0xBB, 0xDD))

# Vertical connector 3->2
add_arrow(slide2, Inches(4.5), Inches(4.95), Inches(4.5), Inches(5.2), GREEN)
add_text_box(slide2, Inches(4.7), Inches(4.95), Inches(0.6), Inches(0.25),
             "DC", font_size=7, font_color=RGBColor(0x99, 0x99, 0x99))

# Layer 2: DC Firewall
y = Inches(5.2)
layer2_details = [
    "Source & Destination Server IP Whitelisting",
    "Only from Specific Server to Specific DMZ Server",
    "IPS, Zoning & Multi-Instance",
    "Unidirectional Flow from DC to DMZ"
]
add_layer_row(slide2, y, 2, GREEN, "DC\nFirewall", layer2_details,
              RGBColor(0x14, 0x24, 0x12), RGBColor(0x36, 0x60, 0x25))

# Vertical connector 2->1
add_arrow(slide2, Inches(4.5), Inches(6.4), Inches(4.5), Inches(6.55), DARK_GREEN)

# Layer 1: ACI
y = Inches(6.55)
layer1_details = [
    "Source Whitelisting - Only from specific EPG to DMZ",
    "Service Graphs, PBR, VRFs",
    "EPGs, ESGs & Multi-tenancy",
    "Unidirectional Flow from EPG to DMZ"
]
add_layer_row(slide2, y, 1, DARK_GREEN, "ACI", layer1_details,
              RGBColor(0x0C, 0x22, 0x15), RGBColor(0x0C, 0x44, 0x2A))

# Left note
add_text_box(slide2, Inches(0.15), Inches(6.7), Inches(1.4), Inches(0.6),
             "Only specific DC servers can communicate unidirectionally with DMZ Servers",
             font_size=6, font_color=RGBColor(0x88, 0xBB, 0xDD))

# ============================================
# SLIDE 3: Traffic Flow Detail
# ============================================
slide3 = prs.slides.add_slide(slide_layout)
bg3 = slide3.background
fill3 = bg3.fill
fill3.solid()
fill3.fore_color.rgb = DARK_BG

# Title
add_rounded_rect(slide3, Inches(3.5), Inches(0.2), Inches(6.3), Inches(0.6),
                 CYAN, text="Traffic Flow - Secure Pull Model", font_size=20, bold=True)

# Flow description
flow_items = [
    ("1. Internet Traffic", "Palo Alto Firewall handles all incoming/outgoing internet traffic with strict IP whitelisting", RED, Inches(1.2)),
    ("2. Proxy Layer", "Web Proxy server processes all requests with URL filtering, malware scanning & HTTPS inspection", ORANGE, Inches(2.4)),
    ("3. DMZ Protection", "DMZ Firewall enforces zone-based policies, allowing only specific DMZ servers to reach proxy", CYAN, Inches(3.6)),
    ("4. DC Communication", "Datacenter Firewall manages internal east-west traffic between DC and DMZ segments", GREEN, Inches(4.8)),
    ("5. ACI Fabric", "Cisco ACI provides micro-segmentation via EPGs, service graphs, and policy-based routing", DARK_GREEN, Inches(6.0)),
]

for title, desc, color, y_pos in flow_items:
    # Color bar
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.0), y_pos, Inches(0.15), Inches(0.9)) if False else \
          slide3.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.0), y_pos, Inches(0.15), Inches(0.9))
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.fill.background()

    # Title
    add_text_box(slide3, Inches(1.4), y_pos, Inches(4.0), Inches(0.4),
                 title, font_size=14, font_color=color, bold=True)
    # Description
    add_text_box(slide3, Inches(1.4), y_pos + Inches(0.4), Inches(10.0), Inches(0.5),
                 desc, font_size=11, font_color=LIGHT_GRAY)

    # Arrow indicator
    arrow_shape = slide3.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(11.5), y_pos + Inches(0.2), Inches(0.6), Inches(0.4))
    arrow_shape.fill.solid()
    arrow_shape.fill.fore_color.rgb = color
    arrow_shape.line.fill.background()

# Key principles box
key_box = add_rounded_rect(slide3, Inches(1.0), Inches(6.8), Inches(11.3), Inches(0.5),
                            RGBColor(0x1B, 0x28, 0x38), border_color=CYAN,
                            text="KEY: All traffic flows are UNIDIRECTIONAL | Each layer enforces IP Whitelisting | Defense-in-Depth Strategy",
                            font_size=10)

# Save the file
prs.save('/projects/sandbox/Secure-Pull/Secure_Pull_Architecture.pptx')
print("PPTX file created successfully!")
